from pptx import Presentation
from typing import Dict, Any
from .base_converter import BaseConverter
from ..utils.medical_terms import extract_medical_terms
from ..utils.table_extractor import extract_tables

class PptxConverter(BaseConverter):
    """Converter for PPTX presentations"""

    def convert(self, file_path: str, **kwargs) -> Dict[str, Any]:
        prs = Presentation(file_path)
        
        # Extract text from slides
        text_content = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text.append(shape.text)
            text_content.append('\n'.join(slide_text))
        
        text = '\n\n=== Новый слайд ===\n\n'.join(text_content)
        
        # Extract tables
        tables = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    tables.append(table_data)
        
        # Extract medical terms
        terms = extract_medical_terms(text)
        
        # Document structure
        structure = {
            'total_slides': len(prs.slides),
            'tables_count': len(tables),
            'has_notes': any(slide.has_notes_slide for slide in prs.slides)
        }
        
        return {
            'text': text,
            'metadata': {
                'slide_count': len(prs.slides),
                'title': prs.core_properties.title
            },
            'tables': tables,
            'terms': terms,
            'structure': structure
        }
    
    def get_supported_formats(self) -> list:
        return ['.pptx', '.ppt']